"""
Script pour générer une présentation PowerPoint professionnelle
pour le projet de prédiction de Customer Churn
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Créer une présentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Définir une couleur de thème cohérente
THEME_COLOR = RGBColor(41, 128, 185)  # Bleu professionnel
ACCENT_COLOR = RGBColor(52, 152, 219)
TEXT_COLOR = RGBColor(44, 62, 80)

def add_title_slide(title, subtitle):
    """Ajouter une diapositive de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Ajouter un rectangle de couleur en arrière-plan
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(2.5))
    background.fill.solid()
    background.fill.fore_color.rgb = THEME_COLOR
    background.line.color.rgb = THEME_COLOR
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = TEXT_COLOR

def add_content_slide(title, content_list):
    """Ajouter une diapositive avec contenu bullet"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = THEME_COLOR
    
    # Ligne décoration sous le titre
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(2), Inches(0))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)
    
    # Contenu
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)

# ======================== DIAPOSITIVE 1: TITRE ========================
add_title_slide(
    "🎯 Prédiction du Customer Churn",
    "Projet End-to-End Machine Learning\nAvril 2026"
)

# ======================== DIAPOSITIVE 2: VUE D'ENSEMBLE ========================
add_content_slide(
    "Vue d'Ensemble du Projet",
    [
        "📌 Objectif: Prédire quels clients vont arrêter leur abonnement",
        "📊 Dataset: 7,043 clients Telco avec 26% de taux de churn",
        "🔍 Enjeu: La rétention coûte moins cher que l'acquisition",
        "✅ Approche: Pipeline ML end-to-end avec déploiement en production",
        "🎁 Livrables: Modèle ML, API FastAPI, MLflow tracking, Docker"
    ]
)

# ======================== DIAPOSITIVE 3: DATASET ========================
add_content_slide(
    "Données et Caractéristiques",
    [
        "📁 Source: Kaggle - Telco Customer Churn Dataset",
        "👥 Taille: 7,043 clients avec 26 features au départ",
        "📈 Déséquilibre: 26% churn vs 74% non-churn (imbalancé)",
        "🔍 Données manquantes: ~11 valeurs manquantes dans TotalCharges",
        "📊 Types de features:",
        "    • Démographiques: Âge, genre, état relationnel",
        "    • Services: Internet, téléphone, streaming, support",
        "    • Contrats: Type, durée, facturation"
    ]
)

# ======================== DIAPOSITIVE 4: EDA - KEY FINDINGS ========================
add_content_slide(
    "Analyse Exploratoire (EDA) - Résultats Clés",
    [
        "✨ Facteurs les plus importants de churn:",
        "    1. Tenure < 12 mois (clients nouveaux)",
        "    2. Contrats mois-à-mois (moins d'engagement)",
        "    3. Paiement par chèque électronique",
        "    4. Absence de contrat d'assistance technique",
        "    5. Services Internet Fiber Optic (problèmes de qualité?)",
        "📊 Distribution: Churn fortement corrélé avec la durée de contrat",
        "💡 Insight business: Cibler les nouveaux clients pour la rétention"
    ]
)

# ======================== DIAPOSITIVE 5: PREPROCESSING ========================
add_content_slide(
    "Prétraitement et Feature Engineering",
    [
        "🔧 Étapes appliquées:",
        "    1. Handling des valeurs manquantes (imputation par médiane)",
        "    2. Encoding des variables catégoriques (LabelEncoder)",
        "    3. Normalisation des features (StandardScaler)",
        "✨ Features d'ingénierie créées:",
        "    • tenure_group: Groupes de durée client",
        "    • avg_monthly_charge: Charge mensuelle moyenne",
        "    • service_count: Nombre de services souscrits",
        "📋 Split: 80/20 en train/test avec stratification"
    ]
)

# ======================== DIAPOSITIVE 6: MODÈLES ========================
add_content_slide(
    "Modèles Entraînés et Performance",
    [
        "🤖 Modèles testés:",
        "    • Logistic Regression (baseline)",
        "    • Random Forest (ensemble)",
        "    • XGBoost (winner) ⭐",
        "📊 Résultats comparatifs (AUC-ROC):",
        "    • Logistic Regression: 0.84",
        "    • Random Forest: 0.87",
        "    • XGBoost: 0.89 ← MEILLEUR",
        "🎯 Métrique retenue: AUC-ROC (meilleure discrimination)",
        "📈 Recall: 61% (capture majorité des churners)"
    ]
)

# ======================== DIAPOSITIVE 7: POURQUOI XGBOOST ========================
add_content_slide(
    "Choix du Meilleur Modèle: XGBoost",
    [
        "✅ AUC-ROC = 0.89 (meilleure discrimination client/non-client)",
        "✅ Recall = 61% vs Precision = 70% (éq. équilibré)",
        "✅ F1-Score meilleur que Logistic Regression et Random Forest",
        "✅ Gestion naturelle du déséquilibre des classes",
        "✅ Un erreur coûtant cher = rater un client qui va partir",
        "💡 Décision business: La rétention est plus importante que",
        "    la sur-ciblage (Recall > Precision)"
    ]
)

# ======================== DIAPOSITIVE 8: FEATURE IMPORTANCE ========================
add_content_slide(
    "Feature Importance (XGBoost)",
    [
        "🔝 Top 5 factors d'importance:",
        "    1. tenure - Durée du contrat (très important!)",
        "    2. Contract - Type de contrat",
        "    3. MonthlyCharges - Frais mensuels",
        "    4. TechSupport - Support technique",
        "    5. InternetService - Type de service Internet",
        "💡 Implication: Les variables d'engagement client sont clés",
        "🎯 Action: Investir dans support et contrats plus longs"
    ]
)

# ======================== DIAPOSITIVE 9: ARCHITECTURE ========================
add_content_slide(
    "Architecture Technique",
    [
        "📦 Structure du projet:",
        "    src/data/ - Chargement et preprocessing",
        "    src/features/ - Feature engineering",
        "    src/models/ - Entraînement et évaluation",
        "    src/visualization/ - Plots et analyses",
        "",
        "🔄 Pipeline ML:",
        "    Data → Preprocessing → Feature Engineering",
        "    → Training → Evaluation → Deployment"
    ]
)

# ======================== DIAPOSITIVE 10: API FASTAPI ========================
add_content_slide(
    "API FastAPI et Déploiement",
    [
        "⚡ FastAPI - Framework web haute performance",
        "📡 Endpoints exposés:",
        "    • POST /predict - Prédiction pour un client",
        "    • GET /health - Vérification de santé",
        "📊 Input: Données client en JSON",
        "📤 Output: Probabilité de churn + niveau de risque",
        "🔗 Documentation: Swagger UI sur http://localhost:8000/docs",
        "✅ Tests: Suite complète de tests unitaires incluse"
    ]
)

# ======================== DIAPOSITIVE 11: MLFLOW TRACKING ========================
add_content_slide(
    "MLflow - Experiment Tracking",
    [
        "📊 MLflow pour versioning et tracking automatique",
        "📝 Chaque entraînement enregistre:",
        "    • Hyperparamètres (learning_rate, max_depth, etc.)",
        "    • Métriques (AUC, Precision, Recall, F1)",
        "    • Artifacts (modèles, scalers, évaluations)",
        "🔄 Comparaison facile entre expériences",
        "📦 Model Registry pour versioning des modèles deployés",
        "🎯 Reproducibilité: Chaque run est traçable"
    ]
)

# ======================== DIAPOSITIVE 12: DOCKER ========================
add_content_slide(
    "Containerisation avec Docker",
    [
        "🐳 Docker Compose pour déploiement complète",
        "📦 Services containerisés:",
        "    • API FastAPI",
        "    • MLflow Server (tracking + model registry)",
        "🚀 One-command deployment: docker-compose up",
        "🔗 Portabilité: Même environment sur n'importe quel serveur",
        "📝 docker-compose.yml configurable pour production",
        "✅ Isolation des dépendances et versions Python"
    ]
)

# ======================== DIAPOSITIVE 13: RÉSULTATS BUSINESS ========================
add_content_slide(
    "Résultats et Impact Business",
    [
        "🎯 Modèle capable de prédire 61% des churners",
        "💰 Impact potentiel:",
        "    • Identifier clients à risque avant départ",
        "    • Personnaliser offres de rétention",
        "    • Réduire coût d'acquisition client",
        "📈 ROI: Coûts de rétention << Coûts de réacquisition",
        "🔄 Amélioration continue: Données de rétention →",
        "    ↺ Entraînement de modèles plus performants"
    ]
)

# ======================== DIAPOSITIVE 14: POINTS FORTS ========================
add_content_slide(
    "Points Forts du Projet",
    [
        "✅ Approche complète: Data → ML → Production",
        "✅ Modèle performant: AUC-ROC 0.89 vs 0.84 baseline",
        "✅ API documentée avec Swagger UI",
        "✅ Tests automatisés pour robustesse",
        "✅ Tracking expériences avec MLflow",
        "✅ Containerisation Docker pour déploiement simple",
        "✅ Code modulaire et réutilisable",
        "✅ Documentation complète (README, GUIDE_DEMO)"
    ]
)

# ======================== DIAPOSITIVE 15: AMÉLIORATIONS FUTURES ========================
add_content_slide(
    "Perspectives d'Améliorations",
    [
        "🔮 Améliorations possibles:",
        "    1. Ajouter features temporelles (saisonnalité)",
        "    2. Ensemble methods pour meilleure robustesse",
        "    3. Feature engineering plus profond",
        "    4. Monitoring et alert en production",
        "    5. A/B testing des stratégies de rétention",
        "    6. Explainability (SHAP) pour iterpétabilité du modèle",
        "    7. Réentraînement automatique avec nouvelles données"
    ]
)

# ======================== DIAPOSITIVE 16: CONCLUSION ========================
add_title_slide(
    "Merci! 🙏",
    "Questions & Discussions\n\ngithub.com/GHASSEN-ABASSI/Customer_Churn_Project"
)

# Sauvegarder la présentation
output_path = r"c:\Users\Ideapad1\Desktop\machine learning projet\Presentation_Customer_Churn.pptx"
prs.save(output_path)
print(f"✅ Présentation créée avec succès: {output_path}")
